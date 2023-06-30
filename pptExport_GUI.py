############## GUI#########################################

from pptx import Presentation
import pandas as pd
import tkinter.ttk as ttk
import tkinter.messagebox as msgbox
from tkinter import * # __all__
import tkinter as tk
from tkinter import filedialog
import threading, time

MAX_ITEM_NO = 12

###간단한 GUI Layout




### GUI Event Functions

sopfile = ''
def open_file():
    global sopfile
    try:
        # sopfile = filedialog.askopenfilename(initialdir=r'C:\Users\ParkGY\DocumentsCFLTYanadoo\DT Academy\SOPMeterial',title="select a file",
        #                                     filetypes =(("Perenstaion","*.pptx"),
        #                                     ("all files","*.*")))
        sopfile = filedialog.askopenfilename(initialdir=r'C:\Users\ParkGY\DocumentsCFLTYanadoo\DT Academy\SOPMeterial',title="select a file",
                                            filetypes =(("Perenstaion","*.pptx"),
                                            ("all files","*.*")))
    except:
        pass




def convert():
    t= threading.Thread(target=threadfunc)
    t.start()

def threadfunc():
    prs = Presentation(sopfile)
    df = pd.DataFrame(columns=['OperationStep','Ref' ,'Man.Item.No','Ser.Item.No' ,'Description','Qty', 'Instruction'])
    
    for i, slide in enumerate(prs.slides):
    #for slide in prs.slides:
        items = GetItemsInSlide(slide) 
        if items : # items list 에 데이터가 있으면 추가 
            op_num = get_opnum(slide.shapes)
            
            for item in items:
                item_dict = {'OperationStep': op_num,'Ref': item[1] ,'Man.Item.No': item[2],'Ser.Item.No': item[3] ,'Description': item[4], 'Qty': item[5], 'Instruction' : item[6]}
                df.loc[len(df)] = item_dict

        progress = (i + 1) / len(prs.slides) * 100
        p_var.set(progress)
        progress_bar.update()

        #Progress Bar 꽉 차면 원위치
        # if progress >= 100:
        #     p_var.set(0)

    if progress == 100:
        window = tk.Tk()

        progress_text = tk.Label(window, text="Convert Completed.")
        progress_text.pack()

        window.mainloop()
                
    f_name = sopfile.split('.')[0] + '.csv' 
    df.to_csv(f_name,encoding='utf-8-sig', index=False, mode='w', header=True)


### Functions

FIRST_ROW = 2
INSTRUCTION_ROW = 10

LEFT_REF_COL = 0
LEFT_ITEM_COL = 2 
LEFT_SER_COL = 4
LEFT_DES_COL = 5
LEFT_QTY_COL = 7
RIGHT_REF_COL = 8
RIGHT_ITEM_COL = 9
RIGHT_SER_COL = 11 
RIGHT_DES_COL = 12
RIGHT_QTY_COL = 14

    

def getTextInRowCol(table, row, col) : 
    text = ''
    for paragraph in table.cell(row,col).text_frame.paragraphs:
        for run in paragraph.runs:
            text += run.text 
    return text

def GetItemInTable(table, idx):        
    
    right_first_item_no = int(MAX_ITEM_NO / 2)
    
    try:
        if idx < right_first_item_no :
            
            op_num = ''
            ref = getTextInRowCol(table, FIRST_ROW + idx, LEFT_REF_COL)
            item_no = getTextInRowCol(table, FIRST_ROW + idx, LEFT_ITEM_COL)
            ser = getTextInRowCol(table, FIRST_ROW + idx, LEFT_SER_COL)
            des = getTextInRowCol(table, FIRST_ROW + idx, LEFT_DES_COL )
            qty = getTextInRowCol(table, FIRST_ROW + idx, LEFT_QTY_COL)
            inst = getTextInRowCol(table,INSTRUCTION_ROW,0 )
            # inst = get_instruction(table.cell)
            # inst =''
            
        else:
            
            op_num = ''
            ref = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_REF_COL)
            item_no = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_ITEM_COL)
            ser = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_SER_COL)
            des = getTextInRowCol(table, FIRST_ROW  + idx - right_first_item_no, RIGHT_DES_COL)
            qty = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_QTY_COL)
                
        item_list = [op_num, ref, item_no, ser, des, qty, inst]         
        return item_list
    except:
        pass
    


def GetItemsInSlide(slide):
    items = []
    
    if slide.shapes[0].has_table:
        table = slide.shapes[0].table

        for idx in range(0, MAX_ITEM_NO):
            try:
                item = GetItemInTable(table, idx)
                if item[5].isnumeric() and item[2]: # Description에 데이터가 있고 qty가 숫자이면 데이터 인정 
                    items.append(item)
                
                elif item[6]:
                    
                    items.append(item)

            except:
                pass
        return items
                

def get_opnum(shapes):
    op_text = ''
    for i, shape in enumerate(shapes):
        if shape.has_text_frame:
            text_in_shape = shape.text[:2]
            if text_in_shape == "OP" and shape.text[2:].isnumeric():
                op_text = shape.text
                break
    if  op_text == '':
        print("coulnd't find OP shape")  
    return op_text

# def get_instruction(table):
#     if table.cell(INSTRUCTION_ROW, 0).text_frame.text_lower() == 'instruction':
#         return table.cell(INSTRUCTION_ROW, 0).text_frame.text
#     return None

# def get_instruction(table):
#     if 'instruction' in table.cell(INSTRUCTION_ROW, 0).text_frame.text.lower():
#         return table.cell(INSTRUCTION_ROW, 0).text_frame.text
#     return None

# def get_instruction(slide):
#     for shape in slide.shapes:
#         if shape.has_table:
#             table = shape.table
#             inst_cell = table.cell(INSTRUCTION_ROW, 0)
#             if 'instruction' in inst_cell.text_frame.text.lower():
                
#                 return inst_cell.text_frame.text
#     return None

# def get_instruction(slide):
#     for shape in slide.shapes:
#         if shape.has_table:
#             table = shape.table
#             for row in table.rows:
#                 if row.cells[0].text_frame.text.lower() == 'instruction':
#                     return row.cells[1].text_frame.text
#     return None


root = Tk()
root.title("PPTX SOP Converter")
root.geometry("200x120")

# Frame 
frame = Frame(root)
frame.pack(fill="x", padx=5, pady=5) # 간격 띄우기

btn_open_file = Button(frame, padx=5, pady=5, text="Open File", command=open_file)
btn_open_file.pack(fill="x",side="top",padx=5, pady=5)

btn_start = Button(frame, padx=5, pady=5, text="Convert", command=convert)
btn_start.pack(fill="x",side="top", padx=5, pady=5)

# 프로그래스 바를 추가
p_var = DoubleVar()
progress_bar = ttk.Progressbar(frame, maximum=100, variable=p_var)
progress_bar.pack(fill="x", padx=5, pady=5)

# Runs
root.mainloop()