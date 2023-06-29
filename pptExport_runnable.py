
## Runable

from pptx import Presentation
import pandas as pd

MAX_ITEM_NO = 12
prs = Presentation('SOPKR6061_01_G3 EUV_ELEC MODULE.pptx')


## Functions

FIRST_ROW = 2
INSTRUCTION_ROW = 10

LEFT_REF_COL = 0 
LEFT_ITEM_COL = 2 
LEFT_SER_COL = 4
LEFT_DES_COL = 5
LEFT_QTY_COL = 7

RIGHT_REF_COL = 8
RIGHT_ITEM_COL = 9
RIGHT_SER_COL = 11 
RIGHT_DES_COL = 12
RIGHT_QTY_COL = 14

    

def getTextInRowCol(table, row, col) : # 테이블 돌면서 텍스트 찾아냄
    text = ''
    for paragraph in table.cell(row,col).text_frame.paragraphs:
        for run in paragraph.runs:
            text += run.text 
    return text

def GetItemInTable(table, idx): # 찾은 텍스트 각 항목에 매핑
    
    right_first_item_no = int(MAX_ITEM_NO / 2)
    
    try:
        if idx < right_first_item_no :
            op_num = ''
            ref = getTextInRowCol(table, FIRST_ROW + idx, LEFT_REF_COL)
            item_no = getTextInRowCol(table, FIRST_ROW + idx, LEFT_ITEM_COL)
            ser = getTextInRowCol(table, FIRST_ROW + idx, LEFT_SER_COL)
            des = getTextInRowCol(table, FIRST_ROW + idx, LEFT_DES_COL )
            qty = getTextInRowCol(table, FIRST_ROW + idx, LEFT_QTY_COL)
            inst = getTextInRowCol(table,INSTRUCTION_ROW,0 )
        else:
            op_num = ''
            ref = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_REF_COL)
            item_no = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_ITEM_COL)
            ser = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_SER_COL)
            des = getTextInRowCol(table, FIRST_ROW  + idx - right_first_item_no, RIGHT_DES_COL)
            qty = getTextInRowCol(table, FIRST_ROW + idx - right_first_item_no, RIGHT_QTY_COL)
                
        item_list = [op_num, ref, item_no, ser, des, qty, inst]         
        return item_list
    except:
        pass
    


def GetItemsInSlide(slide): # 쓰레기값 빼고 아이템들 추가
    items = []
    
    if slide.shapes[0].has_table:
        table = slide.shapes[0].table
        for idx in range(0, MAX_ITEM_NO):
            try:
                item = GetItemInTable(table, idx)
                if item[5].isnumeric() and item[2]: # Description에 데이터가 있고 qty가 숫자이면 데이터 인정 
                    items.append(item)
                    
            except:
                pass
        return items


def get_opnum(shapes):
    op_text = ''
    for i, shape in enumerate(shapes):
        if shape.has_text_frame:
            text_in_shape = shape.text[:2]
            if text_in_shape == "OP":
                op_text = shape.text
    if  op_text == '':
        print("coulnd't find OP shape")  
    return op_text


### 실행코드

df = pd.DataFrame(columns=['OperationStep','Ref', 'Man.Item.No','Ser.Item.No' ,'Description','Qty','Instruction'])

for slide in prs.slides:
    items = GetItemsInSlide(slide) 
    if items : # items list 에 데이터가 있으면 추가 
        op_num = get_opnum(slide.shapes)
        for item in items:
            item_dict = {'OperationStep': op_num, 'Ref': item[1], 'Man.Item.No': item[2], 'Ser.Item.No': item[3], 'Description': item[4], 'Qty': item[5], 'Instruction' : item[6]}
            df.loc[len(df)] = item_dict
            
f_name = 'numbering_sample_2.csv'
df.to_csv(f_name,encoding='utf-8-sig', index=False, mode='w', header=True)